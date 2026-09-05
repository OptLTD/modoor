from __future__ import annotations

from io import BytesIO

import pytest

from sqlalchemy import select

from modoor.core.db import session_scope
from modoor.core.settings import get_settings
from modoor.platform.bootstrap import bootstrap
from modoor.runtime.jobs import Job, run_pending
from modoor.runtime.auth import resolve_ctx
from modules.doc import domain as doc_domain
from modules.doc.extract import extract_bytes, preview_excel, sanitize_pg_text, _clip
from tests.conftest import configure_test_db


@pytest.fixture(autouse=True)
def _env(monkeypatch, tmp_path):
    configure_test_db(
        monkeypatch,
        MODOOR_API_KEY="test-key",
        MODOOR_TENANT="demo",
        MODOOR_CONFIRM_SECRET="secret",
        MODOOR_DOC_STORAGE="local",
        MODOOR_DOC_LOCAL_ROOT=str(tmp_path),
        MODOOR_DOC_OCR="0",
        MODOOR_JOBS_INPROCESS="0",
    )
    bootstrap(get_settings())
    yield
    get_settings.cache_clear()


def test_upload_extracts_on_queue():
    ctx = resolve_ctx(get_settings())
    with session_scope() as session:
        created = doc_domain.create_asset(
            session,
            ctx,
            filename="note.txt",
            data="hello queued extract".encode("utf-8"),
            title="Note",
        )
        assert created["text_status"] == "pending"
        assert created.get("text") in (None, "")
        jobs = list(session.scalars(select(Job).where(Job.kind == "doc.extract")))
        assert len(jobs) == 1
        asset_id = created["id"]

    assert run_pending(limit=4) >= 1
    with session_scope() as session:
        asset = doc_domain.get_asset(session, ctx, asset_id=asset_id)
        assert asset["text_status"] == "ready"
        assert "hello queued extract" in asset["text"]
        assert asset["text_method"] == "text"


def test_provided_text_skips_queue():
    ctx = resolve_ctx(get_settings())
    with session_scope() as session:
        created = doc_domain.create_text_asset(
            session, ctx, title="Memo", text="already here"
        )
        assert created["text_status"] == "ready"
        assert created["text"] == "already here"
        assert session.scalar(select(Job).where(Job.kind == "doc.extract")) is None


def test_extract_docx_xlsx_pptx():
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    doc_buf = BytesIO()
    doc = Document()
    doc.add_paragraph("word body hello")
    doc.save(doc_buf)

    wb = Workbook()
    ws = wb.active
    ws.title = "S1"
    ws["A1"] = "excel body hello"
    xlsx_buf = BytesIO()
    wb.save(xlsx_buf)

    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[-1])
    box = slide.shapes.add_textbox(0, 0, 3000000, 500000)
    box.text_frame.text = "slide body hello"
    ppt_buf = BytesIO()
    ppt.save(ppt_buf)

    assert "word body hello" in extract_bytes("a.docx", doc_buf.getvalue()).text
    assert "excel body hello" in extract_bytes("a.xlsx", xlsx_buf.getvalue()).text
    assert "slide body hello" in extract_bytes("a.pptx", ppt_buf.getvalue()).text


def test_extract_excel_sheet_limits():
    from openpyxl import Workbook

    wb = Workbook()
    small = wb.active
    small.title = "Alpha"
    small["A1"] = "hdr-a"
    small["B1"] = "hdr-b"
    small["A2"] = "keep-me"

    for i in range(1, 7):
        if i == 1:
            continue
        ws = wb.create_sheet(f"S{i}")
        ws["A1"] = f"head-{i}"
        ws["A2"] = f"body-{i}"

    huge = wb.create_sheet("Huge")
    for c in range(1, 51):
        huge.cell(1, c, f"H{c}")
    for r in range(2, 52):
        for c in range(1, 51):
            huge.cell(r, c, f"{r}-{c}")

    buf = BytesIO()
    wb.save(buf)
    text = extract_bytes("cap.xlsx", buf.getvalue()).text
    assert "# sheets: Alpha, S2, S3, S4, S5, S6, Huge" in text
    assert "# Alpha" in text
    assert "hdr-a\thdr-b" in text
    assert "keep-me" in text
    assert "# S5" in text
    assert "head-5" in text
    assert "# S6" not in text
    assert "body-6" not in text
    assert "# Huge" not in text

    huge_only = Workbook()
    hs = huge_only.active
    hs.title = "Big"
    for c in range(1, 51):
        hs.cell(1, c, f"H{c}")
    for r in range(2, 52):
        for c in range(1, 51):
            hs.cell(r, c, f"{r}-{c}")
    buf2 = BytesIO()
    huge_only.save(buf2)
    big = extract_bytes("big.xlsx", buf2.getvalue()).text
    assert "# Big" in big
    header = next(line for line in big.splitlines() if line.startswith("H1\t"))
    assert header.split("\t")[0] == "H1"
    assert header.split("\t")[-1] == "H50"
    assert len(header.split("\t")) == 50
    body_lines = [
        line
        for line in big.splitlines()
        if line and not line.startswith("#") and not line.startswith("H1\t")
    ]
    cells = sum(len(line.split("\t")) for line in body_lines)
    assert cells == 1000
    assert "2-1" in big
    assert "50-50" not in big


def test_preview_excel_pages():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Main"
    ws["A1"] = "col-a"
    ws["B1"] = "col-b"
    for i in range(1, 1006):
        ws.cell(i + 1, 1, f"r{i}")
        ws.cell(i + 1, 2, f"v{i}")
    other = wb.create_sheet("Other")
    other["A1"] = "other-h"
    other["A2"] = "other-row"
    buf = BytesIO()
    wb.save(buf)
    data = buf.getvalue()

    page1 = preview_excel("cap.xlsx", data, sheet=0, page=1)
    assert page1["sheet_names"] == ["Main", "Other"]
    assert page1["sheet"] == 0
    assert page1["header"] == ["col-a", "col-b"]
    assert page1["page"] == 1
    assert page1["page_size"] == 1000
    assert page1["start_row"] == 2
    assert page1["has_prev"] is False
    assert page1["has_next"] is True
    assert len(page1["rows"]) == 1000
    assert page1["rows"][0] == ["r1", "v1"]
    assert page1["rows"][-1] == ["r1000", "v1000"]
    assert page1["row_numbers"][0] == 2
    assert page1["row_numbers"][-1] == 1001

    page2 = preview_excel("cap.xlsx", data, sheet=0, page=2)
    assert page2["page"] == 2
    assert page2["start_row"] == 1002
    assert page2["has_prev"] is True
    assert page2["has_next"] is False
    assert page2["header"] == ["col-a", "col-b"]
    assert len(page2["rows"]) == 5
    assert page2["rows"][0] == ["r1001", "v1001"]
    assert page2["rows"][-1] == ["r1005", "v1005"]

    other_page = preview_excel("cap.xlsx", data, sheet=1, page=1)
    assert other_page["sheet"] == 1
    assert other_page["sheet_name"] == "Other"
    assert other_page["header"] == ["other-h"]
    assert other_page["rows"] == [["other-row"]]
    assert other_page["has_next"] is False


def test_preview_excel_dates():
    from datetime import date, datetime, time

    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "when"
    ws["B1"] = "at"
    ws["C1"] = "clock"
    ws["A2"] = date(2026, 9, 4)
    ws["B2"] = datetime(2026, 9, 4, 15, 30)
    ws["C2"] = time(15, 30)
    ws["A2"].number_format = "YYYY-MM-DD"
    ws["B2"].number_format = "YYYY-MM-DD HH:MM:SS"
    ws["C2"].number_format = "HH:MM:SS"
    buf = BytesIO()
    wb.save(buf)
    page = preview_excel("dates.xlsx", buf.getvalue(), sheet=0, page=1)
    assert page["rows"][0] == ["2026-09-04", "2026-09-04 15:30", "15:30"]


def test_preview_excel_filters():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "status"
    ws["B1"] = "city"
    ws["C1"] = "id"
    statuses = ["open", "closed", "wait"]
    cities = ["SH", "BJ"]
    for i in range(1, 121):
        ws.cell(i + 1, 1, statuses[(i - 1) % 3])
        ws.cell(i + 1, 2, cities[(i - 1) % 2])
        ws.cell(i + 1, 3, f"id-{i}")
    buf = BytesIO()
    wb.save(buf)
    data = buf.getvalue()

    facets = preview_excel("f.xlsx", data, sheet=0, page=1, facets=True)
    names = {item["name"]: item for item in facets["filters"]}
    assert "status" in names
    assert "city" in names
    assert "id" not in names
    assert names["status"]["values"] == [
        {"value": "closed", "count": 40},
        {"value": "open", "count": 40},
        {"value": "wait", "count": 40},
    ]
    assert names["city"]["values"] == [
        {"value": "BJ", "count": 60},
        {"value": "SH", "count": 60},
    ]

    filtered = preview_excel(
        "f.xlsx",
        data,
        sheet=0,
        page=1,
        filters={0: ["open"]},
    )
    assert len(filtered["rows"]) == 40
    assert all(row[0] == "open" for row in filtered["rows"])
    assert filtered["row_numbers"][0] == 2
    assert filtered["row_numbers"][1] == 5
    assert filtered["has_next"] is False

    both = preview_excel(
        "f.xlsx",
        data,
        sheet=0,
        page=1,
        filters={0: ["open", "wait"], 1: ["SH"]},
    )
    assert all(row[0] in {"open", "wait"} and row[1] == "SH" for row in both["rows"])
    assert len(both["rows"]) == 40


def test_sanitize_strips_nul_for_postgres():
    dirty = "KEHOMQ8A\x000013 Sep 3\x00Oct 3"
    assert "\x00" not in sanitize_pg_text(dirty)
    assert _clip(dirty) == "KEHOMQ8A0013 Sep 3Oct 3"
    ctx = resolve_ctx(get_settings())
    with session_scope() as session:
        created = doc_domain.create_text_asset(
            session, ctx, title="nul-invoice", text=dirty
        )
        asset = doc_domain.get_asset(session, ctx, asset_id=created["id"])
        assert "\x00" not in (asset.get("text") or "")
        assert "KEHOMQ8A0013" in (asset.get("text") or "")

