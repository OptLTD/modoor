"""Plain-text extraction: office files + PDF + OCR (images / scanned PDF)."""

from __future__ import annotations

import logging
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from datetime import date, datetime, time
from io import BytesIO
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from modoor.core.settings import get_settings

log = logging.getLogger("modoor.doc.extract")

MAX_TEXT_CHARS = 500_000
EXCEL_MAX_SHEETS = 5
EXCEL_AREA_LIMIT = 1600
EXCEL_BODY_CELLS = 1000
PREVIEW_PAGE_SIZE = 1000
PREVIEW_MAX_COLS = 40
PREVIEW_FILTER_MAX_UNIQ = 100
_TEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".json", ".log", ".html", ".htm", ".xml"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".tif", ".tiff", ".bmp"}
_OCR_ENGINE: Any = None
_OCR_FAILED = False


@dataclass
class ExtractResult:
    text: str
    method: str
    error: str = ""


def extract_bytes(filename: str, data: bytes, mime: str = "") -> ExtractResult:
    ext = Path(filename or "").suffix.lower()
    mime_l = (mime or "").lower()
    try:
        if ext in _TEXT_EXTS or mime_l.startswith("text/") or mime_l in (
            "application/json",
            "application/xml",
            "application/javascript",
        ):
            return ExtractResult(_decode_text(data), "text")
        if ext in _IMAGE_EXTS or mime_l.startswith("image/"):
            return _extract_image(data)
        if ext == ".pdf" or mime_l == "application/pdf":
            return _extract_pdf(data)
        if ext == ".docx" or mime_l.endswith("wordprocessingml.document"):
            return _extract_docx(data)
        if ext == ".doc" or mime_l == "application/msword":
            return _extract_doc(data)
        if ext in {".xlsx", ".xlsm"} or "spreadsheetml.sheet" in mime_l:
            return _extract_xlsx(data)
        if ext == ".xls" or mime_l == "application/vnd.ms-excel":
            return _extract_xls(data)
        if ext == ".pptx" or "presentationml.presentation" in mime_l:
            return _extract_pptx(data)
        if ext == ".ppt" or mime_l == "application/vnd.ms-powerpoint":
            return _extract_ppt(data)
        return ExtractResult("", "none")
    except Exception as exc:  # noqa: BLE001
        log.exception("extract failed for %s", filename)
        return ExtractResult("", "error", str(exc)[:2000])


def sanitize_pg_text(text: str) -> str:
    """Postgres text/varchar cannot store NUL (0x00); PDF extract often emits them."""
    return (text or "").replace("\x00", "")


def _clip(text: str) -> str:
    return sanitize_pg_text(text).strip()[:MAX_TEXT_CHARS]


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return sanitize_pg_text(data.decode(enc))[:MAX_TEXT_CHARS]
        except UnicodeDecodeError:
            continue
    return sanitize_pg_text(data.decode("utf-8", errors="replace"))[:MAX_TEXT_CHARS]


def _ocr_enabled() -> bool:
    return bool(get_settings().modoor_doc_ocr)


def _ocr_engine():
    global _OCR_ENGINE, _OCR_FAILED
    if _OCR_FAILED or not _ocr_enabled():
        return None
    if _OCR_ENGINE is not None:
        return _OCR_ENGINE
    try:
        from rapidocr_onnxruntime import RapidOCR

        _OCR_ENGINE = RapidOCR()
        return _OCR_ENGINE
    except Exception:  # noqa: BLE001
        log.warning("RapidOCR unavailable", exc_info=True)
        _OCR_FAILED = True
        return None


def _ocr_image_bytes(data: bytes) -> str:
    engine = _ocr_engine()
    if engine is None:
        return ""
    try:
        import numpy as np
        from PIL import Image

        img = Image.open(BytesIO(data))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        w, h = img.size
        if w < 8 or h < 8:
            return ""
        if max(w, h) > 2000:
            scale = 2000 / max(w, h)
            img = img.resize((max(int(w * scale), 1), max(int(h * scale), 1)))
        arr = np.asarray(img.convert("RGB"))
        result, _elapse = engine(arr)
        if not result:
            return ""
        lines: list[str] = []
        for item in result:
            if isinstance(item, (list, tuple)) and len(item) >= 2:
                lines.append(str(item[1]))
        return "\n".join(lines).strip()
    except Exception:  # noqa: BLE001
        log.exception("ocr image failed")
        return ""


def _extract_image(data: bytes) -> ExtractResult:
    if not _ocr_enabled():
        return ExtractResult("", "ocr_disabled")
    text = _ocr_image_bytes(data)
    if not text and _OCR_FAILED:
        return ExtractResult("", "ocr_unavailable", "RapidOCR is not available")
    return ExtractResult(_clip(text), "ocr")


def _extract_pdf(data: bytes) -> ExtractResult:
    methods: list[str] = []
    page_texts: list[str] = []
    n_pages = 0
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        n_pages = len(reader.pages)
        for page in reader.pages:
            page_texts.append((page.extract_text() or "").strip())
        methods.append("pdf")
    except Exception:  # noqa: BLE001
        log.debug("pypdf failed", exc_info=True)
        page_texts = []

    max_ocr = int(get_settings().modoor_doc_ocr_max_pages)
    need_ocr_idx = [
        i
        for i, t in enumerate(page_texts)
        if len(t) < 20 and i < max_ocr
    ]
    if not page_texts:
        need_ocr_idx = list(range(min(max(n_pages, 1), max_ocr)))

    if need_ocr_idx and _ocr_enabled():
        ocr_pages = _ocr_pdf_pages(data, need_ocr_idx)
        if ocr_pages:
            methods.append("ocr")
            if not page_texts:
                page_texts = [""] * (max(need_ocr_idx) + 1)
            for i, text in ocr_pages.items():
                while len(page_texts) <= i:
                    page_texts.append("")
                if len(text.strip()) > len(page_texts[i].strip()):
                    page_texts[i] = text

    joined = "\n\n".join(t for t in page_texts if t.strip())
    method = "+".join(methods) if methods else "pdf"
    return ExtractResult(_clip(joined), method)


def _ocr_pdf_pages(data: bytes, indexes: list[int]) -> dict[int, str]:
    engine = _ocr_engine()
    if engine is None:
        return {}
    try:
        import pypdfium2 as pdfium
    except Exception:  # noqa: BLE001
        log.warning("pypdfium2 unavailable", exc_info=True)
        return {}
    out: dict[int, str] = {}
    try:
        pdf = pdfium.PdfDocument(data)
        for i in indexes:
            if i < 0 or i >= len(pdf):
                continue
            page = pdf[i]
            bitmap = page.render(scale=2)
            pil = bitmap.to_pil()
            buf = BytesIO()
            pil.save(buf, format="PNG")
            text = _ocr_image_bytes(buf.getvalue())
            if text:
                out[i] = text
            page.close()
        pdf.close()
    except Exception:  # noqa: BLE001
        log.exception("pdf ocr render failed")
    return out


def _xml_t_texts(xml: bytes) -> list[str]:
    root = ET.fromstring(xml)
    return [(node.text or "") for node in root.iter() if node.tag.endswith("}t") and node.text]


def _extract_docx(data: bytes) -> ExtractResult:
    parts: list[str] = []
    try:
        from docx import Document

        doc = Document(BytesIO(data))
        parts.extend(p.text for p in doc.paragraphs if p.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append("\t".join(cells))
        if _ocr_enabled():
            try:
                for rel in doc.part.rels.values():
                    if "image" in getattr(rel, "reltype", ""):
                        blob = rel.target_part.blob
                        ocr = _ocr_image_bytes(blob)
                        if ocr:
                            parts.append(ocr)
            except Exception:  # noqa: BLE001
                log.debug("docx image ocr skipped", exc_info=True)
        if parts:
            return ExtractResult(_clip("\n".join(parts)), "docx")
    except Exception:  # noqa: BLE001
        log.debug("python-docx failed, trying zip xml", exc_info=True)
    try:
        with ZipFile(BytesIO(data)) as zf:
            xml = zf.read("word/document.xml")
        texts = _xml_t_texts(xml)
        return ExtractResult(_clip("\n".join(texts)), "docx")
    except Exception as exc:  # noqa: BLE001
        return ExtractResult("", "docx", str(exc)[:2000])


def _cell_str(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        value = _format_datetime(value)
    elif isinstance(value, date):
        value = value.isoformat()
    elif isinstance(value, time):
        value = _format_time(value)
    return str(value).replace("\t", " ").replace("\r", " ").replace("\n", " ").strip()


def _format_time(value: time) -> str:
    if value.second or value.microsecond:
        return value.replace(microsecond=0).strftime("%H:%M:%S")
    return value.strftime("%H:%M")


def _format_datetime(value: datetime) -> str:
    value = value.replace(tzinfo=None, microsecond=0)
    d = value.date()
    t = value.time()
    # Excel time-only serials surface as 1899-12-30 + clock
    if (d.year, d.month, d.day) == (1899, 12, 30):
        return _format_time(t)
    if t == time(0, 0, 0):
        return d.isoformat()
    return f"{d.isoformat()} {_format_time(t)}"


def _trim_row(row: list[Any] | tuple[Any, ...]) -> list[str]:
    vals = [_cell_str(c) for c in row]
    while vals and vals[-1] == "":
        vals.pop()
    return vals


def _pad_row(row: list[str], ncols: int) -> list[str]:
    if len(row) >= ncols:
        return row[:ncols]
    return row + [""] * (ncols - len(row))


def _take_body_cells(body: list[list[str]], ncols: int, limit: int) -> list[list[str]]:
    out: list[list[str]] = []
    taken = 0
    for row in body:
        if taken >= limit:
            break
        padded = _pad_row(row, ncols)
        n = min(ncols, limit - taken)
        out.append(padded[:n])
        taken += n
    return out


def _format_one_sheet(title: str, grid: list[list[str]]) -> list[str]:
    lines = [f"# {title}"]
    if not grid:
        return lines
    ncols = max((len(r) for r in grid), default=0)
    header = _pad_row(grid[0], ncols)
    lines.append("\t".join(header))
    body = grid[1:]
    area = len(grid) * max(ncols, 1)
    if area > EXCEL_AREA_LIMIT:
        body = _take_body_cells(body, ncols, EXCEL_BODY_CELLS)
        for row in body:
            lines.append("\t".join(row))
    else:
        for row in body:
            lines.append("\t".join(_pad_row(row, ncols)))
    return lines


def _format_workbook(all_names: list[str], extracted: list[tuple[str, list[list[str]]]]) -> str:
    parts: list[str] = []
    if all_names:
        parts.append("# sheets: " + ", ".join(all_names))
    for title, grid in extracted:
        parts.extend(_format_one_sheet(title, grid))
    return "\n".join(parts)


def _grid_from_rows(rows) -> list[list[str]]:
    it = iter(rows)
    try:
        first = next(it)
    except StopIteration:
        return []
    header = _trim_row(list(first))
    ncols = max(len(header), 1)
    grid: list[list[str]] = [header]
    for raw in it:
        row = _trim_row(list(raw))
        if len(row) > ncols:
            ncols = len(row)
        grid.append(row)
        if len(grid) * ncols > EXCEL_AREA_LIMIT:
            body = _take_body_cells(grid[1:], ncols, EXCEL_BODY_CELLS)
            return [header, *body]
    return grid


def _preview_row(row: list[Any] | tuple[Any, ...], max_cols: int) -> list[str]:
    vals = [_cell_str(c) for c in list(row)[:max_cols]]
    while vals and vals[-1] == "":
        vals.pop()
    return vals


def _xls_row(sheet: Any, row: int, take_c: int) -> list[Any]:
    import xlrd
    from xlrd.xldate import xldate_as_datetime

    if take_c <= 0:
        return []
    datemode = sheet.book.datemode
    out: list[Any] = []
    for c in range(take_c):
        cell = sheet.cell(row, c)
        if cell.ctype == xlrd.XL_CELL_DATE:
            try:
                out.append(xldate_as_datetime(cell.value, datemode))
                continue
            except (ValueError, OverflowError):
                pass
        out.append(cell.value)
    return out


def preview_excel(
    filename: str,
    data: bytes,
    mime: str = "",
    *,
    sheet: int = 0,
    page: int = 1,
    filters: dict[Any, Any] | None = None,
    facets: bool = False,
) -> dict[str, Any]:
    """One sheet, one page (1000 data rows). Header is always included."""
    ext = Path(filename or "").suffix.lower()
    mime_l = (mime or "").lower()
    page = max(int(page or 1), 1)
    sheet = max(int(sheet or 0), 0)
    parsed = _coerce_filters(filters)
    if ext == ".xls" or mime_l == "application/vnd.ms-excel":
        return _preview_xls(
            data, sheet=sheet, page=page, filters=parsed, facets=facets
        )
    return _preview_xlsx(
        data, sheet=sheet, page=page, filters=parsed, facets=facets
    )


def _coerce_filters(filters: dict[Any, Any] | None) -> dict[int, set[str]]:
    out: dict[int, set[str]] = {}
    if not filters:
        return out
    for key, raw in filters.items():
        try:
            col = int(key)
        except (TypeError, ValueError):
            continue
        if col < 0 or col >= PREVIEW_MAX_COLS:
            continue
        if not isinstance(raw, (list, tuple, set)):
            continue
        allowed = {"" if v is None else str(v) for v in raw}
        if allowed:
            out[col] = allowed
    return out


def _match_range(page: int) -> tuple[int, int]:
    """1-based match ordinal among rows that pass filters (data rows only)."""
    lo = 1 + (page - 1) * PREVIEW_PAGE_SIZE
    hi = page * PREVIEW_PAGE_SIZE
    return lo, hi


def _cell_at(row: list[str], col: int) -> str:
    return row[col] if 0 <= col < len(row) else ""


def _row_matches(row: list[str], filters: dict[int, set[str]]) -> bool:
    for col, allowed in filters.items():
        if _cell_at(row, col) not in allowed:
            return False
    return True


def _facet_payload(header: list[str], acc: "_FacetAcc") -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for col in sorted(acc.counts):
        tallies = acc.counts[col]
        if not tallies or set(tallies) == {""} or len(tallies) >= PREVIEW_FILTER_MAX_UNIQ:
            continue
        name = _cell_at(header, col) or f"col{col + 1}"
        ordered = sorted(tallies.items(), key=lambda kv: (kv[0] == "", kv[0].lower(), kv[0]))
        out.append(
            {
                "col": col,
                "name": name,
                "values": [{"value": value, "count": count} for value, count in ordered],
            }
        )
    return out


class _FacetAcc:
    def __init__(self) -> None:
        self.counts: dict[int, dict[str, int]] = {}
        self.dead: set[int] = set()

    @property
    def needed(self) -> bool:
        return len(self.dead) < PREVIEW_MAX_COLS

    def add(self, row: list[str]) -> None:
        if not self.needed:
            return
        width = min(max(len(row), 0), PREVIEW_MAX_COLS)
        for col in list(self.counts):
            if col >= width:
                self._put(col, "")
        for col in range(width):
            if col in self.dead:
                continue
            self._put(col, _cell_at(row, col))

    def _put(self, col: int, value: str) -> None:
        bucket = self.counts.setdefault(col, {})
        if value not in bucket and len(bucket) + 1 >= PREVIEW_FILTER_MAX_UNIQ:
            self.dead.add(col)
            del self.counts[col]
            return
        bucket[value] = bucket.get(value, 0) + 1


def _collect_preview_page(
    row_iter,
    *,
    page: int,
    filters: dict[int, set[str]],
    facets: bool,
) -> tuple[list[str], list[list[str]], list[int], bool, list[dict[str, Any]] | None]:
    match_lo, match_hi = _match_range(page)
    header: list[str] = []
    rows: list[list[str]] = []
    row_numbers: list[int] = []
    match_n = 0
    has_next = False
    acc = _FacetAcc() if facets else None
    for excel_i, raw in enumerate(row_iter, start=1):
        if excel_i == 1:
            header = _preview_row(raw, PREVIEW_MAX_COLS)
            continue
        if not filters and acc is None:
            match_n += 1
            if match_n < match_lo:
                continue
            if match_n > match_hi:
                has_next = True
                break
            rows.append(_preview_row(raw, PREVIEW_MAX_COLS))
            row_numbers.append(excel_i)
            continue
        cells = _preview_row(raw, PREVIEW_MAX_COLS)
        if acc is not None:
            acc.add(cells)
        if filters and not _row_matches(cells, filters):
            continue
        match_n += 1
        if match_n < match_lo:
            continue
        if match_n > match_hi:
            has_next = True
            if acc is None or not acc.needed:
                break
            continue
        rows.append(cells)
        row_numbers.append(excel_i)
    return header, rows, row_numbers, has_next, (_facet_payload(header, acc) if acc else None)


def _preview_payload(
    *,
    names: list[str],
    sheet: int,
    header: list[str],
    rows: list[list[str]],
    row_numbers: list[int],
    page: int,
    has_next: bool,
    filters: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    if not names:
        names = ["Sheet1"]
    if sheet >= len(names):
        sheet = 0
    start = row_numbers[0] if row_numbers else 2
    payload: dict[str, Any] = {
        "kind": "xlsx",
        "sheet_names": names,
        "sheet": sheet,
        "sheet_name": names[sheet],
        "header": header,
        "rows": rows,
        "row_numbers": row_numbers,
        "page": page,
        "page_size": PREVIEW_PAGE_SIZE,
        "start_row": start,
        "has_prev": page > 1,
        "has_next": has_next,
    }
    if filters is not None:
        payload["filters"] = filters
    return payload


def _preview_xlsx(
    data: bytes,
    *,
    sheet: int,
    page: int,
    filters: dict[int, set[str]],
    facets: bool,
) -> dict[str, Any]:
    import openpyxl

    wb = openpyxl.load_workbook(BytesIO(data), data_only=True, read_only=True)
    try:
        names = [s.title for s in wb.worksheets]
        if not names:
            return _preview_payload(
                names=[],
                sheet=0,
                header=[],
                rows=[],
                row_numbers=[],
                page=page,
                has_next=False,
                filters=[] if facets else None,
            )
        if sheet >= len(names):
            sheet = 0
        ws = wb.worksheets[sheet]
        header, rows, row_numbers, has_next, facet_rows = _collect_preview_page(
            ws.iter_rows(max_col=PREVIEW_MAX_COLS, values_only=True),
            page=page,
            filters=filters,
            facets=facets,
        )
        return _preview_payload(
            names=names,
            sheet=sheet,
            header=header,
            rows=rows,
            row_numbers=row_numbers,
            page=page,
            has_next=has_next,
            filters=facet_rows,
        )
    finally:
        wb.close()


def _preview_xls(
    data: bytes,
    *,
    sheet: int,
    page: int,
    filters: dict[int, set[str]],
    facets: bool,
) -> dict[str, Any]:
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    names = [s.name for s in book.sheets()]
    if not names:
        return _preview_payload(
            names=[],
            sheet=0,
            header=[],
            rows=[],
            row_numbers=[],
            page=page,
            has_next=False,
            filters=[] if facets else None,
        )
    if sheet >= len(names):
        sheet = 0
    ws = book.sheet_by_index(sheet)
    take_c = min(max(int(ws.ncols or 0), 0), PREVIEW_MAX_COLS)

    def _iter_xls():
        for r in range(int(ws.nrows or 0)):
            yield _xls_row(ws, r, take_c)

    header, rows, row_numbers, has_next, facet_rows = _collect_preview_page(
        _iter_xls(),
        page=page,
        filters=filters,
        facets=facets,
    )
    return _preview_payload(
        names=names,
        sheet=sheet,
        header=header,
        rows=rows,
        row_numbers=row_numbers,
        page=page,
        has_next=has_next,
        filters=facet_rows,
    )


def _extract_xlsx(data: bytes) -> ExtractResult:
    try:
        import openpyxl

        wb = openpyxl.load_workbook(BytesIO(data), data_only=True, read_only=True)
        names = [sheet.title for sheet in wb.worksheets]
        extracted: list[tuple[str, list[list[str]]]] = []
        for sheet in wb.worksheets[:EXCEL_MAX_SHEETS]:
            extracted.append((sheet.title, _grid_from_rows(sheet.iter_rows(values_only=True))))
        wb.close()
        return ExtractResult(_clip(_format_workbook(names, extracted)), "xlsx")
    except Exception:  # noqa: BLE001
        log.debug("openpyxl failed, trying zip xml", exc_info=True)
    return _xlsx_from_zip(data)


def _xlsx_from_zip(data: bytes) -> ExtractResult:
    try:
        with ZipFile(BytesIO(data)) as zf:
            names = _xlsx_sheet_names(zf)
            shared: list[str] = []
            if "xl/sharedStrings.xml" in zf.namelist():
                ss = ET.fromstring(zf.read("xl/sharedStrings.xml"))
                for si in ss:
                    shared.append(
                        "".join((t.text or "") for t in si.iter() if t.tag.endswith("}t"))
                    )
            sheet_files = [
                n
                for n in zf.namelist()
                if n.startswith("xl/worksheets/sheet") and n.endswith(".xml")
            ]
            sheet_files.sort()
            extracted: list[tuple[str, list[list[str]]]] = []
            for i, name in enumerate(sheet_files[:EXCEL_MAX_SHEETS]):
                title = names[i] if i < len(names) else Path(name).stem
                extracted.append((title, _xlsx_sheet_grid(zf.read(name), shared)))
            if not names:
                names = [title for title, _ in extracted]
            return ExtractResult(_clip(_format_workbook(names, extracted)), "xlsx")
    except Exception as exc:  # noqa: BLE001
        return ExtractResult("", "xlsx", str(exc)[:2000])


def _xlsx_sheet_names(zf: ZipFile) -> list[str]:
    if "xl/workbook.xml" not in zf.namelist():
        return []
    root = ET.fromstring(zf.read("xl/workbook.xml"))
    names: list[str] = []
    for node in root.iter():
        if node.tag.endswith("}sheet") or node.tag == "sheet":
            name = node.get("name")
            if name:
                names.append(name)
    return names


def _xlsx_sheet_grid(xml: bytes, shared: list[str]) -> list[list[str]]:
    root = ET.fromstring(xml)
    rows: dict[int, dict[int, str]] = {}
    for c in root.iter():
        if not c.tag.endswith("}c"):
            continue
        ref = c.get("r") or ""
        col_s = "".join(ch for ch in ref if ch.isalpha())
        row_s = "".join(ch for ch in ref if ch.isdigit())
        if not col_s or not row_s:
            continue
        col = 0
        for ch in col_s.upper():
            col = col * 26 + (ord(ch) - 64)
        row_i = int(row_s)
        v = next((ch for ch in list(c) if ch.tag.endswith("}v")), None)
        if v is None or v.text is None:
            continue
        if c.get("t") == "s":
            try:
                text = shared[int(v.text)]
            except (ValueError, IndexError):
                text = v.text
        else:
            text = v.text
        rows.setdefault(row_i, {})[col] = _cell_str(text)
    if not rows:
        return []
    grid: list[list[str]] = []
    for r in sorted(rows):
        cols = rows[r]
        width = max(cols)
        grid.append([cols.get(c, "") for c in range(1, width + 1)])
        ncols = max((len(row) for row in grid), default=1)
        if len(grid) * ncols > EXCEL_AREA_LIMIT:
            header = grid[0]
            body = _take_body_cells(grid[1:], ncols, EXCEL_BODY_CELLS)
            return [header, *body]
    return grid


def _extract_xls(data: bytes) -> ExtractResult:
    try:
        import xlrd

        book = xlrd.open_workbook(file_contents=data)
        names = [sheet.name for sheet in book.sheets()]
        extracted: list[tuple[str, list[list[str]]]] = []
        for sheet in book.sheets()[:EXCEL_MAX_SHEETS]:
            ncols = max(int(sheet.ncols or 0), 1)
            nrows = int(sheet.nrows or 0)
            if nrows * ncols > EXCEL_AREA_LIMIT:
                extra = max(1, (EXCEL_BODY_CELLS + ncols - 1) // ncols)
                take = min(nrows, 1 + extra)
            else:
                take = nrows
            rows = [_xls_row(sheet, r, ncols) for r in range(take)]
            extracted.append((sheet.name, _grid_from_rows(rows)))
        return ExtractResult(_clip(_format_workbook(names, extracted)), "xls")
    except Exception as exc:  # noqa: BLE001
        return ExtractResult("", "xls", str(exc)[:2000])


def _extract_pptx(data: bytes) -> ExtractResult:
    parts: list[str] = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(BytesIO(data))
        for i, slide in enumerate(prs.slides, 1):
            slide_bits: list[str] = [f"# slide {i}"]
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = (shape.text or "").strip()
                    if text:
                        slide_bits.append(text)
                if _ocr_enabled() and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        ocr = _ocr_image_bytes(blob)
                        if ocr:
                            slide_bits.append(ocr)
                    except Exception:  # noqa: BLE001
                        log.debug("pptx picture ocr skipped", exc_info=True)
            notes = getattr(slide, "has_notes_slide", False)
            if notes:
                try:
                    note = slide.notes_slide.notes_text_frame.text
                    if note.strip():
                        slide_bits.append(note.strip())
                except Exception:  # noqa: BLE001
                    pass
            if len(slide_bits) > 1:
                parts.append("\n".join(slide_bits))
        if parts:
            return ExtractResult(_clip("\n\n".join(parts)), "pptx")
    except Exception:  # noqa: BLE001
        log.debug("python-pptx failed, trying zip xml", exc_info=True)
    try:
        with ZipFile(BytesIO(data)) as zf:
            for name in sorted(zf.namelist()):
                if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
                    continue
                parts.extend(_xml_t_texts(zf.read(name)))
        return ExtractResult(_clip("\n".join(parts)), "pptx")
    except Exception as exc:  # noqa: BLE001
        return ExtractResult("", "pptx", str(exc)[:2000])


def _run_bin(cmd: list[str], data: bytes, suffix: str) -> str:
    if not shutil.which(cmd[0]):
        return ""
    with tempfile.NamedTemporaryFile(suffix=suffix) as tmp:
        tmp.write(data)
        tmp.flush()
        try:
            proc = subprocess.run(
                [*cmd, tmp.name],
                capture_output=True,
                timeout=60,
                check=False,
            )
        except (OSError, subprocess.TimeoutExpired):
            return ""
    if proc.returncode != 0:
        return ""
    return proc.stdout.decode("utf-8", errors="replace").strip()


def _ole_strings(data: bytes) -> str:
    try:
        import olefile
    except Exception:  # noqa: BLE001
        return ""
    try:
        ole = olefile.OleFileIO(BytesIO(data))
    except Exception:  # noqa: BLE001
        return ""
    chunks: list[str] = []
    try:
        for stream in ole.listdir():
            try:
                raw = ole.openstream(stream).read()
            except Exception:  # noqa: BLE001
                continue
            text = _binary_printable(raw)
            if text:
                chunks.append(text)
    finally:
        ole.close()
    return max(chunks, key=len) if chunks else ""


def _binary_printable(raw: bytes) -> str:
    candidates: list[str] = []
    for enc in ("utf-16le", "utf-8", "latin-1"):
        try:
            s = raw.decode(enc, errors="ignore")
        except Exception:  # noqa: BLE001
            continue
        cleaned = "".join(ch if ch.isprintable() or ch in "\n\t" else " " for ch in s)
        cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
        tokens = re.findall(r"[\w\u4e00-\u9fff]{3,}", cleaned)
        if len(tokens) >= 8:
            candidates.append(cleaned)
    return max(candidates, key=len) if candidates else ""


def _extract_doc(data: bytes) -> ExtractResult:
    text = _run_bin(["antiword"], data, ".doc") or _run_bin(["catdoc"], data, ".doc")
    if text:
        return ExtractResult(_clip(text), "doc")
    ole = _ole_strings(data)
    if ole:
        return ExtractResult(_clip(ole), "doc")
    return ExtractResult("", "doc", "legacy .doc: install antiword/catdoc for better extraction")


def _extract_ppt(data: bytes) -> ExtractResult:
    text = _run_bin(["catppt"], data, ".ppt")
    if text:
        return ExtractResult(_clip(text), "ppt")
    ole = _ole_strings(data)
    if ole:
        return ExtractResult(_clip(ole), "ppt")
    return ExtractResult("", "ppt", "legacy .ppt: install catppt for better extraction")
